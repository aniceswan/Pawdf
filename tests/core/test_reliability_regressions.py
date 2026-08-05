"""Regression tests for data-preserving reliability fixes."""

from __future__ import annotations

from pathlib import Path
from types import SimpleNamespace

import pytest


@pytest.fixture
def complete_form_pdf(tmp_path):
    from reportlab.lib.pagesizes import letter
    from reportlab.pdfgen import canvas

    path = tmp_path / "complete-form.pdf"
    output = canvas.Canvas(str(path), pagesize=letter)
    output.drawString(72, 740, "Application")
    form = output.acroForm
    form.textfield(name="fullname", x=72, y=690, width=220, height=24, value="")
    form.checkbox(name="agree", x=72, y=650, checked=False)
    form.radio(name="plan", value="basic", x=72, y=610, selected=True)
    form.radio(name="plan", value="pro", x=110, y=610, selected=False)
    form.choice(
        name="country",
        value="ID",
        options=["ID", "US"],
        x=72,
        y=560,
        width=100,
        height=24,
    )
    output.showPage()
    output.save()
    return path


def test_radio_group_is_one_field_with_all_options(complete_form_pdf):
    from pawdf.core.forms import list_fields

    radios = [field for field in list_fields(complete_form_pdf) if field.name == "plan"]

    assert len(radios) == 1
    assert {option.removeprefix("/") for option in radios[0].options} >= {"basic", "pro"}


def test_radio_selection_turns_on_only_the_chosen_widget(complete_form_pdf, tmp_path):
    import pikepdf

    from pawdf.core.forms import fill_form, list_fields

    out = fill_form(complete_form_pdf, tmp_path / "radio.pdf", {"plan": "/pro"})
    radio = next(field for field in list_fields(out) if field.name == "plan")
    assert radio.value == "/pro"

    states = []
    with pikepdf.open(out) as pdf:
        for page in pdf.pages:
            for annot in page.get("/Annots", []):
                if annot.get("/Subtype") != "/Widget":
                    continue
                parent = annot.get("/Parent")
                if parent is not None and str(parent.get("/T", "")) == "plan":
                    states.append(str(annot.get("/AS", "/Off")))

    assert states.count("/pro") == 1
    assert states.count("/basic") == 0


def test_choice_field_lists_and_accepts_real_options(complete_form_pdf, tmp_path):
    from pawdf.core.forms import fill_form, list_fields

    choice = next(field for field in list_fields(complete_form_pdf) if field.name == "country")
    assert choice.kind == "choice"
    assert set(choice.options) >= {"ID", "US"}

    out = fill_form(complete_form_pdf, tmp_path / "choice.pdf", {"country": "US"})
    updated = next(field for field in list_fields(out) if field.name == "country")
    assert updated.value == "US"


def test_flattening_keeps_the_filled_value_visible(complete_form_pdf, tmp_path):
    import pypdfium2 as pdfium

    from pawdf.core.forms import fill_form, list_fields

    out = fill_form(
        complete_form_pdf,
        tmp_path / "flat.pdf",
        {"fullname": "Grace Hopper", "plan": "/pro", "country": "US"},
        flatten=True,
    )
    assert list_fields(out) == []

    document = pdfium.PdfDocument(out)
    try:
        page = document[0]
        try:
            textpage = page.get_textpage()
            try:
                visible_text = textpage.get_text_range()
            finally:
                textpage.close()
        finally:
            page.close()
    finally:
        document.close()

    assert "Grace Hopper" in visible_text


def test_ocr_partial_failure_preserves_page_count_and_order(tmp_path, monkeypatch):
    import pikepdf

    import pawdf.core.ocr as ocr

    source = tmp_path / "source.pdf"
    with pikepdf.new() as pdf:
        for width in (200, 220, 240):
            pdf.add_blank_page(page_size=(width, 300))
        pdf.save(source)

    fake_binary = tmp_path / "tesseract"
    fake_binary.write_text("", encoding="utf-8")
    monkeypatch.setattr(ocr, "find_tesseract", lambda: fake_binary)

    def fake_render(document, index, dpi, into):  # noqa: ANN001, ARG001
        image = into / f"page-{index:05d}.png"
        image.write_bytes(b"fake")
        return image

    def fake_run(args, **kwargs):  # noqa: ANN001, ARG001
        stem = Path(args[2])
        index = int(stem.name.rsplit("-", 1)[1])
        if index == 1:
            return SimpleNamespace(returncode=1, stderr="page failed", stdout="")
        with pikepdf.new() as part:
            part.add_blank_page(page_size=(200 + 20 * index, 300))
            part.save(stem.with_suffix(".pdf"))
        return SimpleNamespace(returncode=0, stderr="", stdout="")

    monkeypatch.setattr(ocr, "_render_page", fake_render)
    monkeypatch.setattr(ocr.subprocess, "run", fake_run)

    result = ocr.ocr_pdf(source, tmp_path / "out.pdf")

    assert result.failed_pages == [2]
    assert result.pages_processed == 2
    with pikepdf.open(result.output_path) as output:
        assert len(output.pages) == 3
        widths = [round(float(page.MediaBox[2])) for page in output.pages]
        assert widths == [200, 220, 240]


def test_ocr_all_page_failures_still_return_the_original_pages(tmp_path, monkeypatch):
    import pikepdf

    import pawdf.core.ocr as ocr

    source = tmp_path / "source.pdf"
    with pikepdf.new() as pdf:
        pdf.add_blank_page(page_size=(200, 300))
        pdf.save(source)

    fake_binary = tmp_path / "tesseract"
    fake_binary.write_text("", encoding="utf-8")
    monkeypatch.setattr(ocr, "find_tesseract", lambda: fake_binary)
    monkeypatch.setattr(
        ocr,
        "_render_page",
        lambda document, index, dpi, into: into / "fake.png",
    )
    monkeypatch.setattr(
        ocr.subprocess,
        "run",
        lambda *args, **kwargs: SimpleNamespace(
            returncode=1,
            stderr="generic recognition failure",
            stdout="",
        ),
    )

    result = ocr.ocr_pdf(source, tmp_path / "out.pdf")

    assert result.pages_processed == 0
    assert result.failed_pages == [1]
    with pikepdf.open(result.output_path) as output:
        assert len(output.pages) == 1


def test_powerpoint_reports_unsupported_shapes(tmp_path):
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches

    from pawdf.core.pptx_to_pdf import pptx_to_pdf_detailed

    deck = Presentation()
    slide = deck.slides.add_slide(deck.slide_layouts[6])
    slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(1),
        Inches(1),
        Inches(2),
        Inches(1),
    )
    source = tmp_path / "unsupported-shape.pptx"
    deck.save(source)

    result = pptx_to_pdf_detailed(source, tmp_path / "out.pdf")

    assert result.output_path.is_file()
    assert result.shapes_skipped >= 1
    assert result.warnings


def test_ghostscript_hash_mismatch_is_deleted_and_never_accepted(tmp_path):
    from pawdf.core.compress.ghostscript import GhostscriptSetupError, _verify_installer

    installer = tmp_path / "ghostscript.exe"
    installer.write_bytes(b"not the pinned Artifex installer")

    with pytest.raises(GhostscriptSetupError, match="SHA-256"):
        _verify_installer(installer)

    assert not installer.exists()
