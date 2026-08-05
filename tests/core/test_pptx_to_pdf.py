import pytest

from pawdf.core._shared import ConversionError, page_count
from pawdf.core.pptx_to_pdf import pptx_to_pdf


def _text_of(path, index=0):
    import pypdfium2 as pdfium

    doc = pdfium.PdfDocument(path)
    try:
        textpage = doc[index].get_textpage()
        try:
            return textpage.get_text_range()
        finally:
            textpage.close()
    finally:
        doc.close()


@pytest.fixture
def deck(tmp_path):
    from pptx import Presentation

    presentation = Presentation()

    title = presentation.slides.add_slide(presentation.slide_layouts[0])
    title.shapes.title.text = "Quarterly Review"
    title.placeholders[1].text = "Finance team"

    bullets = presentation.slides.add_slide(presentation.slide_layouts[1])
    bullets.shapes.title.text = "Highlights"
    frame = bullets.placeholders[1].text_frame
    frame.text = "Revenue up 12%"
    frame.add_paragraph().text = "Costs flat"

    path = tmp_path / "deck.pptx"
    presentation.save(path)
    return path


def test_one_page_per_slide(deck, tmp_path):
    out = pptx_to_pdf(deck, tmp_path / "out.pdf")
    assert page_count(out) == 2


def test_slide_text_reaches_the_page(deck, tmp_path):
    out = pptx_to_pdf(deck, tmp_path / "out.pdf")

    assert "Quarterly Review" in _text_of(out, 0)
    assert "Revenue up 12%" in _text_of(out, 1)
    assert "Costs flat" in _text_of(out, 1)


def test_page_size_follows_the_deck(tmp_path):
    """A 16:9 deck should not come out on A4 with bars down the side."""
    import pypdfium2 as pdfium
    from pptx import Presentation
    from pptx.util import Emu

    presentation = Presentation()
    presentation.slide_width = Emu(12192000)  # 13.333in, the 16:9 standard
    presentation.slide_height = Emu(6858000)  # 7.5in
    presentation.slides.add_slide(presentation.slide_layouts[6])
    path = tmp_path / "wide.pptx"
    presentation.save(path)

    out = pptx_to_pdf(path, tmp_path / "out.pdf")

    doc = pdfium.PdfDocument(out)
    try:
        width, height = doc[0].get_size()
    finally:
        doc.close()
    assert width / height == pytest.approx(16 / 9, abs=0.02)


def test_pictures_are_embedded(tmp_path):
    import pikepdf
    from PIL import Image
    from pptx import Presentation
    from pptx.util import Inches

    image = tmp_path / "pic.png"
    Image.new("RGB", (200, 120), (30, 90, 160)).save(image)

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.add_picture(str(image), Inches(1), Inches(2), width=Inches(3))
    path = tmp_path / "pic.pptx"
    presentation.save(path)

    out = pptx_to_pdf(path, tmp_path / "out.pdf")

    with pikepdf.open(out) as pdf:
        assert pdf.pages[0].get_images()


def test_an_empty_deck_still_produces_a_readable_pdf(tmp_path):
    """A canvas with no pages writes a file no reader will open."""
    from pptx import Presentation

    path = tmp_path / "empty.pptx"
    Presentation().save(path)

    out = pptx_to_pdf(path, tmp_path / "out.pdf")
    assert page_count(out) == 1


def test_unreadable_input_raises(tmp_path):
    bad = tmp_path / "bad.pptx"
    bad.write_bytes(b"not a deck")

    with pytest.raises(ConversionError):
        pptx_to_pdf(bad, tmp_path / "out.pdf")


def test_missing_input_raises(tmp_path):
    with pytest.raises(FileNotFoundError):
        pptx_to_pdf(tmp_path / "nope.pptx", tmp_path / "out.pdf")
