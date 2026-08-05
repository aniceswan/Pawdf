"""PowerPoint (.pptx) to PDF with explicit fidelity reporting."""

from __future__ import annotations

import io
import warnings
from dataclasses import dataclass
from pathlib import Path

from pawdf.core._shared import ConversionError, ensure_exists, ensure_parent_dir

__all__ = [
    "PptxConversionResult",
    "PresentationFidelityWarning",
    "pptx_to_pdf",
    "pptx_to_pdf_detailed",
]

_EMU_PER_POINT = 12700
_DEFAULT_FONT_PT = 18.0
_MIN_FONT_PT = 6.0


class PresentationFidelityWarning(UserWarning):
    """Warns that a deck contained content Pawdf could not fully reproduce."""


@dataclass
class PptxConversionResult:
    output_path: Path
    warnings: list[str]
    shapes_rendered: int
    shapes_skipped: int


def _pt(emu: int | None, fallback: float = 0.0) -> float:
    return (emu / _EMU_PER_POINT) if emu is not None else fallback


def _shape_font(paragraph, run=None) -> tuple[str, float, tuple[float, float, float]]:  # noqa: ANN001
    size = None
    bold = False
    italic = False
    colour = (0.1, 0.1, 0.1)

    sources = [paragraph.font]
    if run is not None:
        sources.insert(0, run.font)

    for source in sources:
        if size is None and source.size is not None:
            size = source.size.pt
        if source.bold:
            bold = True
        if source.italic:
            italic = True
        try:
            if source.color is not None and source.color.rgb is not None:
                rgb = str(source.color.rgb)
                colour = tuple(int(rgb[i : i + 2], 16) / 255 for i in (0, 2, 4))
        except (AttributeError, ValueError, TypeError):
            pass

    name = "Helvetica"
    if bold and italic:
        name = "Helvetica-BoldOblique"
    elif bold:
        name = "Helvetica-Bold"
    elif italic:
        name = "Helvetica-Oblique"

    return name, max(size or _DEFAULT_FONT_PT, _MIN_FONT_PT), colour


def _draw_text_frame(canvas, shape, page_height: float) -> None:  # noqa: ANN001
    from pptx.enum.text import PP_ALIGN

    left = _pt(shape.left)
    top = _pt(shape.top)
    width = _pt(shape.width, 100.0)
    y = page_height - top

    for paragraph in shape.text_frame.paragraphs:
        text = "".join(run.text for run in paragraph.runs) or paragraph.text
        if not text.strip():
            y -= _DEFAULT_FONT_PT * 0.6
            continue

        first_run = paragraph.runs[0] if paragraph.runs else None
        font_name, size, colour = _shape_font(paragraph, first_run)
        y -= size * 1.2
        canvas.setFont(font_name, size)
        canvas.setFillColorRGB(*colour)

        if paragraph.alignment == PP_ALIGN.CENTER:
            canvas.drawCentredString(left + width / 2, y, text)
        elif paragraph.alignment == PP_ALIGN.RIGHT:
            canvas.drawRightString(left + width, y, text)
        else:
            canvas.drawString(left + 14 * (paragraph.level or 0), y, text)


def _draw_picture(canvas, shape, page_height: float) -> None:  # noqa: ANN001
    from reportlab.lib.utils import ImageReader

    image = shape.image
    reader = ImageReader(io.BytesIO(image.blob))
    width = _pt(shape.width, 100.0)
    height = _pt(shape.height, 100.0)
    canvas.drawImage(
        reader,
        _pt(shape.left),
        page_height - _pt(shape.top) - height,
        width=width,
        height=height,
        mask="auto",
        preserveAspectRatio=True,
    )


def _type_name(shape) -> str:  # noqa: ANN001
    raw = getattr(shape, "shape_type", "unknown")
    name = getattr(raw, "name", None)
    return str(name or raw).replace("_", " ").lower()


def pptx_to_pdf_detailed(
    input_path: str | Path,
    out_path: str | Path,
) -> PptxConversionResult:
    """Convert a deck and report every skipped or partially rendered shape."""
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from reportlab.pdfgen.canvas import Canvas

    source = ensure_exists(input_path)
    out_path = ensure_parent_dir(out_path)

    try:
        deck = Presentation(str(source))
    except Exception as exc:  # noqa: BLE001 - normalized for callers
        raise ConversionError(f"Could not read '{source}': {exc}") from exc

    page_width = _pt(deck.slide_width, 720.0)
    page_height = _pt(deck.slide_height, 540.0)
    fidelity_warnings: list[str] = []
    rendered = 0
    skipped = 0

    try:
        canvas = Canvas(str(out_path), pagesize=(page_width, page_height))
        text_only_types = {
            MSO_SHAPE_TYPE.TEXT_BOX,
            MSO_SHAPE_TYPE.PLACEHOLDER,
        }

        for slide_index, slide in enumerate(deck.slides, start=1):
            for shape_index, shape in enumerate(slide.shapes, start=1):
                shape_type = shape.shape_type
                label = f"Slide {slide_index}, shape {shape_index} ({_type_name(shape)})"
                try:
                    if shape_type == MSO_SHAPE_TYPE.PICTURE:
                        _draw_picture(canvas, shape, page_height)
                        rendered += 1
                    elif shape.has_text_frame:
                        text = shape.text_frame.text.strip()
                        if text:
                            _draw_text_frame(canvas, shape, page_height)
                            rendered += 1
                        if shape_type not in text_only_types:
                            fidelity_warnings.append(
                                f"{label}: text was retained, but the shape's "
                                "fill, outline or effects were not reproduced."
                            )
                            if not text:
                                skipped += 1
                    else:
                        skipped += 1
                        fidelity_warnings.append(
                            f"{label}: this shape type is unsupported and was skipped."
                        )
                except Exception as exc:  # noqa: BLE001 - reported, never silent
                    skipped += 1
                    fidelity_warnings.append(f"{label}: could not render it ({exc}).")
            canvas.showPage()

        if not deck.slides:
            canvas.showPage()
        canvas.save()
    except Exception as exc:  # noqa: BLE001 - normalized for callers
        out_path.unlink(missing_ok=True)
        raise ConversionError(f"Could not convert the presentation to PDF: {exc}") from exc

    return PptxConversionResult(
        output_path=out_path,
        warnings=fidelity_warnings,
        shapes_rendered=rendered,
        shapes_skipped=skipped,
    )


def pptx_to_pdf(input_path: str | Path, out_path: str | Path) -> Path:
    """Compatibility API that warns instead of hiding fidelity loss."""
    result = pptx_to_pdf_detailed(input_path, out_path)
    for message in result.warnings:
        warnings.warn(message, PresentationFidelityWarning, stacklevel=2)
    return result.output_path
